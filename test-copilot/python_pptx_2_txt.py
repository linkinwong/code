from pptx import Presentation
import sys

def pptx_to_text(filename):
    prs = Presentation(filename)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
    return text

if __name__ == "__main__":
    print(pptx_to_text(sys.argv[1]))